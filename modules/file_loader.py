import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import streamlit as st

class FileLoader:
    def __init__(self):
        pass
    
    def load_pdf(self, file):
        """Extract text from PDF"""
        try:
            reader = PdfReader(file)
            text = ""
            total_pages = len(reader.pages)
            
            # Limit to first 50 pages for speed
            pages_to_read = min(total_pages, 50)
            
            for i in range(pages_to_read):
                page_text = reader.pages[i].extract_text()
                if page_text:
                    text += page_text + "\n"
            
            if total_pages > 50:
                st.info(f"📄 Processed first 50 of {total_pages} pages for faster performance")
            
            return text.strip()
        except Exception as e:
            st.error(f"Error reading PDF: {e}")
            return ""
    
    def load_docx(self, file):
        """Extract text from DOCX"""
        try:
            doc = Document(file)
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            return text.strip()
        except Exception as e:
            st.error(f"Error reading DOCX: {e}")
            return ""
    
    def load_pptx(self, file):
        """Extract text from PPTX"""
        try:
            prs = Presentation(file)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            st.error(f"Error reading PPTX: {e}")
            return ""
    
    def load_image(self, file):
        """Process image - return placeholder text for Gemini Vision"""
        try:
            from PIL import Image
            img = Image.open(file)
            
            # Store image info instead of OCR
            st.info("📸 Image uploaded! AI will analyze the image content directly.")
            
            # Return a placeholder that signals we have an image
            return f"[IMAGE_CONTENT: {file.name}]"
            
        except Exception as e:
            st.error(f"Error loading image: {e}")
            return ""
    
    def load_txt(self, file):
        """Load plain text file"""
        try:
            text = file.read().decode("utf-8")
            return text.strip()
        except Exception as e:
            st.error(f"Error reading text file: {e}")
            return ""
    
    def load_multiple_files(self, files):
        """Load multiple files and combine text"""
        combined_text = ""
        
        for file in files:
            st.info(f"Processing: {file.name}")
            
            if file.type == "application/pdf":
                text = self.load_pdf(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text = self.load_docx(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = self.load_pptx(file)
            elif file.type.startswith("image/"):
                text = self.load_image(file)
            elif file.type == "text/plain":
                text = self.load_txt(file)
            else:
                st.warning(f"Unsupported file type: {file.type}")
                continue
            
            if text:
                combined_text += f"\n\n--- Content from {file.name} ---\n\n{text}"
        
        return combined_text.strip()